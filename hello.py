from io import StringIO, BytesIO
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from os import path, system
from markdown_strings import header, esc_format, image as md_img
from pdf2image import convert_from_path
from imgkit import from_string
from markdown import markdown
from base64 import b64encode
import ollama
import asyncio


def process_powerpoint(file_path: str):
    prs = Presentation(file_path)
    output = StringIO()

    output.write(header(f"Current file: {file_path}\n", 1))

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                output.write(esc_format(shape.text, esc=True))
                output.write("\n")

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                output.write(md_img(b64encode(shape.image.blob), shape.image.ext))

    return from_string(markdown(output.getvalue()))


def process_pdf(file_path: str):
    encoded_images = []
    images = convert_from_path(file_path)

    for image in images:
        buffered = BytesIO()
        image.save(buffered, format="JPEG")
        encoded_images.append(b64encode(buffered.getvalue()))

    return encoded_images


async def image_to_latex(
    client: ollama.AsyncClient, image: str, prompt: str, output_tex
) -> str:
    response = await client.chat(
        model="llama3.2-vision",
        messages=[
            {
                "role": "user",
                "content": prompt,
                "images": [image],
            }
        ],
    )

    print(response)
    output_tex.write(response["messages"][0]["content"])


async def main():
    num_files: int = 1  # int(input("Enter number of files: "))

    images = []

    for _ in range(num_files):
        file_path: str = input("Enter path to file: ")

        if not path.isfile(file_path):
            print(f"File {file_path} does not exist")
            continue

        match path.splitext(file_path)[1]:
            case ".pptx":
                images.extend(process_powerpoint(file_path))
            case ".pdf":
                images.extend(process_pdf(file_path))
            case _:
                print(f"File {path.basename(file_path)} is not supported")

    with open("output.tex", "w") as output_tex:
        with open("header.tex", "r") as f:
            output_tex.write(f.read())

        async with asyncio.TaskGroup() as tg:
            client = ollama.AsyncClient()

            with open("prompt.txt", "r") as f:
                prompt = f.read()
                for image in images:
                    tg.create_task(image_to_latex(client, image, prompt, output_tex))

        output_tex.write(r"\end{multicols}\end{document}")

    system("pdflatex output.tex")

    print("All done")


# file -> /Users/benjamingluck/Downloads/35.GUI.pdf
if __name__ == "__main__":
    asyncio.run(main())
